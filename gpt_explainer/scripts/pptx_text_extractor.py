from typing import Generator
from pptx import Presentation
from pptx.exc import PackageNotFoundError

class PptxParser:
    """
    A class used to parse PPTX files and extract text from slides.
    """

    def __init__(self, pptx_file_path: str):
        """
        Initialize PptxParser with a PPTX file.

        Parameters:
            file_path (str): The path to the PPTX file.
        Raises:
            FileNotFoundError: If the file is not found.
        """
        try:
            self.presentation = Presentation(pptx_file_path)
        except (PackageNotFoundError):
            raise FileNotFoundError(f"File not found: {pptx_file_path}")
        

    def extract_text_from_single_slide(self, slide) -> str:
        """
        Extracts and returns all text from a single slide formatted as a string.

        Parameters:
            slide: A slide object from a PPTX file.

        Returns:
            str: A string representing the text from the slide.
        """
        return '\n'.join([' '.join(shape.text.split()) for shape in slide.shapes if hasattr(shape, "text")])
    
    def generate_text_from_slides(self) -> Generator[str, None, None]:
        """
        Extracts and yields all text from the slides in a PPTX file.

        Yields:
            str: A string representing the text from a slide. If an error occurs, the error message is yielded.        
        """
        for i, slide in enumerate(self.presentation.slides):
            try:
                slide_text = self.extract_text_from_single_slide(slide)
                if slide_text:
                    yield slide_text
            except Exception as e:
                yield f"Error processing slide {i+1}: {str(e)}"

